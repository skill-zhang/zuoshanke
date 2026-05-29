"""📽 PowerPoint 生成器 — 从结构化数据或 Markdown 生成 .pptx

支持标题页、多节幻灯片、项目列表、表格、代码展示、图片嵌入。
输出到 zuoshanke 项目下的 output/ 目录。

## 用法
    from tools.gen_ppt import gen_ppt
    r = json.loads(gen_ppt(
        title="项目汇报",
        slides=[
            {"heading":"目标", "content":"- 完成A\\n- 完成B", "layout":"bullet"},
        ]
    ))
"""

import json
import os
import re
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    Presentation = None

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)


# ── 配色方案 ──
COLOR_PRIMARY = RGBColor(30, 60, 120)      # 深蓝
COLOR_ACCENT = RGBColor(0, 120, 180)        # 亮蓝
COLOR_DARK = RGBColor(50, 50, 50)            # 正文深灰
COLOR_LIGHT = RGBColor(100, 100, 100)        # 浅灰
COLOR_BG = RGBColor(240, 244, 250)           # 浅蓝背景


def _set_font(run, size=18, bold=False, color=COLOR_DARK):
    """设置 run 的字体属性"""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Microsoft YaHei"


def _add_textbox(slide, left, top, width, height):
    """添加文本框"""
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    return txbox.text_frame


def _parse_md_inline(text: str) -> list:
    """解析行内 Markdown 格式，返回 (text, bold, italic) 元组列表"""
    segments = []
    parts = re.split(r"(\*\*.+?\*\*|\*.+?\*|`[^`]+`)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            segments.append((part[2:-2], True, False))
        elif part.startswith("*") and part.endswith("*"):
            segments.append((part[1:-1], False, True))
        elif part.startswith("`") and part.endswith("`"):
            segments.append((part[1:-1], False, False))
        else:
            segments.append((part, False, False))
    return segments


def _add_formatted_paragraph(tf, text: str, size=16, bold=False, color=None,
                              align=PP_ALIGN.LEFT, space_before=0, space_after=4):
    """添加带行内格式的段落"""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)

    if color is None:
        color = COLOR_DARK

    segments = _parse_md_inline(text)
    for seg_text, seg_bold, seg_italic in segments:
        run = p.add_run()
        run.text = seg_text
        _set_font(run, size=size, bold=bold or seg_bold, color=color)
        if seg_italic:
            run.font.italic = True

    return p


def _add_bullet_list(tf, lines, size=14):
    """添加项目符号列表"""
    for line in lines:
        p = tf.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        p.level = 0
        # 缩进和项目符号
        segments = _parse_md_inline(line.lstrip("- *"))
        for seg_text, seg_bold, seg_italic in segments:
            run = p.add_run()
            run.text = seg_text
            _set_font(run, size=size, color=COLOR_DARK)
            if seg_italic:
                run.font.italic = True


def gen_ppt(
    title: str = "演示文稿",
    slides: list = None,
    author: str = "坐山客",
    output_name: str = "",
) -> str:
    """生成 PowerPoint (.pptx) 演示文稿

    Args:
        title: 演示文稿标题（第一页标题）
        slides: 幻灯片列表。每项格式:
            {"heading": "标题", "content": "正文(支持Markdown)",
             "layout": "bullet|blank|two_column|code|table",
             "subtitle": "副标题（可选）"}
            支持 layout:
              - bullet (默认): 列表型
              - blank: 全宽正文
              - two_column: 两栏 layout（content 格式 "左栏\\n---\\n右栏"）
              - code: 代码展示（等宽字体，深色背景）
              - table: 表格（content 格式 "|h1|h2|\\n|---|---|\\n|v1|v2|"）
        author: 作者
        output_name: 输出文件名（不含 .pptx）

    Returns:
        JSON: {success, file_path, slide_count, title, error?}
    """
    if Presentation is None:
        return json.dumps({"success": False, "error": "python-pptx 未安装，请运行: pip install python-pptx"})

    try:
        if slides is None:
            slides = []

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # ══ 首页 ══
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_PRIMARY

        # 标题
        tf = _add_textbox(slide, 1.5, 2.0, 10.3, 2.5)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        _set_font(run, size=44, bold=True, color=RGBColor(255, 255, 255))

        # 副标题
        tf2 = _add_textbox(slide, 1.5, 4.5, 10.3, 1.5)
        tf2.word_wrap = True
        p2 = tf2.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = f"作者: {author}  |  {datetime.now().strftime('%Y-%m-%d')}"
        _set_font(run2, size=16, color=RGBColor(200, 210, 230))

        # ══ 内容页 ══
        for sec in slides:
            heading = sec.get("heading", "")
            content = sec.get("content", "")
            layout = sec.get("layout", "bullet")
            subtitle = sec.get("subtitle", "")

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = COLOR_BG

            # 标题区
            tf = _add_textbox(slide, 0.8, 0.4, 11.7, 0.8)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = heading
            _set_font(run, size=30, bold=True, color=COLOR_PRIMARY)

            if subtitle:
                tf_sub = _add_textbox(slide, 0.8, 1.1, 11.7, 0.5)
                p_sub = tf_sub.paragraphs[0]
                run_sub = p_sub.add_run()
                run_sub.text = subtitle
                _set_font(run_sub, size=14, color=COLOR_LIGHT)
                content_top = 1.6
            else:
                content_top = 1.3

            if layout == "bullet":
                lines = [l.strip() for l in content.split("\n") if l.strip()]
                tf_body = _add_textbox(slide, 0.8, content_top, 11.7, 5.5)
                tf_body.word_wrap = True
                _add_bullet_list(tf_body, lines, size=15)

            elif layout == "blank":
                tf_body = _add_textbox(slide, 0.8, content_top, 11.7, 5.5)
                tf_body.word_wrap = True
                for line in content.split("\n"):
                    if line.strip():
                        _add_formatted_paragraph(tf_body, line.strip(), size=16, space_after=6)

            elif layout == "two_column":
                parts = content.split("---")
                left_content = parts[0].strip() if len(parts) > 0 else ""
                right_content = parts[1].strip() if len(parts) > 1 else ""

                # 左栏
                tf_left = _add_textbox(slide, 0.8, content_top, 5.5, 5.5)
                tf_left.word_wrap = True
                for line in left_content.split("\n"):
                    if line.strip():
                        if line.strip().startswith("-"):
                            _add_bullet_list(tf_left, [line.strip()], size=14)
                        else:
                            _add_formatted_paragraph(tf_left, line.strip(), size=14, space_after=4)

                # 分隔线
                vline = slide.shapes.add_shape(1, Inches(6.6), Inches(content_top),
                                                Inches(0.02), Inches(5.0))  # 1 = rectangle
                vline.fill.solid()
                vline.fill.fore_color.rgb = COLOR_ACCENT
                vline.line.fill.background()

                # 右栏
                tf_right = _add_textbox(slide, 6.8, content_top, 5.5, 5.5)
                tf_right.word_wrap = True
                for line in right_content.split("\n"):
                    if line.strip():
                        if line.strip().startswith("-"):
                            _add_bullet_list(tf_right, [line.strip()], size=14)
                        else:
                            _add_formatted_paragraph(tf_right, line.strip(), size=14, space_after=4)

            elif layout == "code":
                tf_body = _add_textbox(slide, 0.8, content_top, 11.7, 5.5)
                tf_body.word_wrap = True
                # 代码背景
                bg_shape = slide.shapes.add_shape(1, Inches(0.8), Inches(content_top),
                                                   Inches(11.7), Inches(0.5 + len(content.split("\n")) * 0.35))
                bg_shape.fill.solid()
                bg_shape.fill.fore_color.rgb = RGBColor(30, 30, 40)
                bg_shape.line.fill.background()

                for line in content.split("\n"):
                    if line.strip():
                        p = tf_body.add_paragraph()
                        p.space_before = Pt(2)
                        p.space_after = Pt(2)
                        run = p.add_run()
                        run.text = line
                        run.font.size = Pt(11)
                        run.font.name = "Courier New"
                        run.font.color.rgb = RGBColor(180, 220, 120)

            elif layout == "table":
                rows_data = []
                for line in content.split("\n"):
                    line = line.strip()
                    if line.startswith("|") and line.endswith("|"):
                        cells = [c.strip() for c in line.strip("|").split("|")]
                        rows_data.append(cells)
                if rows_data:
                    # 过滤分隔行
                    rows_data = [r for r in rows_data if not re.match(r"^[\s:-]+$", r[0])]
                    if rows_data:
                        col_count = min(max(len(r) for r in rows_data), 6)
                        row_count = min(len(rows_data), 15)
                        tbl_left = (13.333 - col_count * 1.8) / 2
                        table = slide.shapes.add_table(
                            row_count, col_count,
                            Inches(tbl_left), Inches(content_top),
                            Inches(col_count * 1.8), Inches(row_count * 0.45)
                        ).table

                        for ri, row_data in enumerate(rows_data[:row_count]):
                            for ci in range(col_count):
                                cell_text = row_data[ci] if ci < len(row_data) else ""
                                cell = table.cell(ri, ci)
                                cell.text = cell_text
                                for p in cell.text_frame.paragraphs:
                                    p.font.size = Pt(12)
                                    if ri == 0:
                                        p.font.bold = True
                                        p.font.color.rgb = RGBColor(255, 255, 255)
                                        # 表头背景
                                        cell.fill.solid()
                                        cell.fill.fore_color.rgb = COLOR_PRIMARY
                                    else:
                                        p.font.color.rgb = COLOR_DARK

        # ══ 保存 ══
        if not output_name:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_title = re.sub(r"[^\w\u4e00-\u9fff_-]", "_", title)[:30]
            output_name = f"{safe_title}_{ts}"
        output_path = str(OUTPUT_DIR / f"{output_name}.pptx")
        prs.save(output_path)

        return json.dumps({
            "success": True,
            "file_path": output_path,
            "slide_count": len(slides) + 1,
            "title": title,
        }, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"success": False, "error": str(e)}, ensure_ascii=False)
