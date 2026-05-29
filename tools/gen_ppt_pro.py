"""
📽 通用 PPT 生成器（专业版）— 多风格、多模板、自动配图

支持 6 种风格配色 + 7 种布局模板 + 自动搜索配图
注册为坐山客工具后，可通过 gen_ppt_pro 调用

## 用法
    from tools.gen_ppt_pro import gen_ppt_pro
    r = json.loads(gen_ppt_pro(
        title="报告标题",
        slides=[...],
        style="tech",  # tech/fashion/food/medical/culture/study
    ))
"""

import json
import os
import re
import urllib.request
import urllib.parse
import random
import string
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    Presentation = None

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)
TEMP_DIR = Path(__file__).resolve().parent.parent / "data" / "temp_ppt_images"
TEMP_DIR.mkdir(parents=True, exist_ok=True)

# ============================================================
# 6 套配色方案
# ============================================================
STYLES = {
    "tech": {
        "name": "科技风",
        "scene": "AI/技术/产品对比",
        "bg_dark": RGBColor(0x0A, 0x0E, 0x1A),
        "bg_dark2": RGBColor(0x0F, 0x18, 0x2E),
        "card_bg": RGBColor(0x14, 0x20, 0x3C),
        "card_bg2": RGBColor(0x1A, 0x2A, 0x4A),
        "primary": RGBColor(0x00, 0xD4, 0xE8),      # 霓虹青
        "secondary": RGBColor(0x00, 0xE0, 0x7A),     # 霓虹绿
        "accent": RGBColor(0xF0, 0xC0, 0x40),        # 金色
        "highlight": RGBColor(0xF0, 0x7A, 0x3A),     # 橙
        "purple": RGBColor(0xA0, 0x5A, 0xE0),        # 紫
        "blue": RGBColor(0x30, 0x90, 0xF0),          # 蓝
        "red": RGBColor(0xE0, 0x40, 0x50),           # 红
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light": RGBColor(0xD8, 0xE0, 0xEE),
        "gray": RGBColor(0x88, 0x99, 0xAA),
        "table_h": RGBColor(0x00, 0x99, 0xBB),
        "table_r1": RGBColor(0x14, 0x22, 0x40),
        "table_r2": RGBColor(0x1A, 0x2C, 0x50),
        "bg_overlay": "99",
        "bg_img_keywords": "technology digital network abstract dark blue",
    },
    "fashion": {
        "name": "时尚风",
        "scene": "时装/美妆/设计",
        "bg_dark": RGBColor(0x1A, 0x10, 0x18),
        "bg_dark2": RGBColor(0x25, 0x18, 0x22),
        "card_bg": RGBColor(0x30, 0x1E, 0x2A),
        "card_bg2": RGBColor(0x3A, 0x28, 0x34),
        "primary": RGBColor(0xE8, 0xA0, 0xB0),      # 玫瑰金
        "secondary": RGBColor(0xB0, 0xA0, 0xD0),     # 烟灰紫
        "accent": RGBColor(0xF0, 0xC0, 0x80),        # 香槟金
        "highlight": RGBColor(0xD0, 0x60, 0x80),     # 玫红
        "purple": RGBColor(0xC0, 0x80, 0xD0),        # 淡紫
        "blue": RGBColor(0x80, 0xA0, 0xD0),          # 灰蓝
        "red": RGBColor(0xD0, 0x50, 0x60),           # 红
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light": RGBColor(0xE8, 0xE0, 0xE8),
        "gray": RGBColor(0xAA, 0x99, 0xA8),
        "table_h": RGBColor(0xC0, 0x70, 0x90),
        "table_r1": RGBColor(0x30, 0x1E, 0x2A),
        "table_r2": RGBColor(0x3A, 0x28, 0x34),
        "bg_overlay": "99",
        "bg_img_keywords": "fashion elegant luxury gradient pastel",
    },
    "food": {
        "name": "美食风",
        "scene": "餐饮/菜谱/食品",
        "bg_dark": RGBColor(0x1A, 0x12, 0x0A),
        "bg_dark2": RGBColor(0x25, 0x1A, 0x0F),
        "card_bg": RGBColor(0x30, 0x22, 0x14),
        "card_bg2": RGBColor(0x3A, 0x2C, 0x1A),
        "primary": RGBColor(0xF0, 0x7A, 0x3A),      # 暖橙
        "secondary": RGBColor(0x7A, 0xA0, 0x3A),     # 橄榄绿
        "accent": RGBColor(0xF0, 0xC0, 0x40),        # 金黄
        "highlight": RGBColor(0xD0, 0x50, 0x30),     # 辣椒红
        "purple": RGBColor(0xA0, 0x60, 0x80),        # 紫
        "blue": RGBColor(0x50, 0x80, 0xA0),          # 蓝
        "red": RGBColor(0xE0, 0x40, 0x30),           # 红
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light": RGBColor(0xE8, 0xE0, 0xD0),
        "gray": RGBColor(0xAA, 0x99, 0x80),
        "table_h": RGBColor(0xD0, 0x70, 0x30),
        "table_r1": RGBColor(0x30, 0x22, 0x14),
        "table_r2": RGBColor(0x3A, 0x2C, 0x1A),
        "bg_overlay": "99",
        "bg_img_keywords": "food cuisine delicious warm organic",
    },
    "medical": {
        "name": "医疗风",
        "scene": "医疗/健康/医药",
        "bg_dark": RGBColor(0x0A, 0x14, 0x12),
        "bg_dark2": RGBColor(0x0F, 0x1E, 0x1A),
        "card_bg": RGBColor(0x14, 0x28, 0x24),
        "card_bg2": RGBColor(0x1A, 0x32, 0x2E),
        "primary": RGBColor(0x00, 0xC8, 0xA0),      # 薄荷绿
        "secondary": RGBColor(0x30, 0x90, 0xF0),     # 医用蓝
        "accent": RGBColor(0x40, 0xE0, 0xC0),        # 青绿
        "highlight": RGBColor(0xE0, 0x70, 0x50),     # 暖橙
        "purple": RGBColor(0x80, 0x60, 0xC0),        # 紫
        "blue": RGBColor(0x30, 0x90, 0xF0),          # 蓝
        "red": RGBColor(0xE0, 0x40, 0x50),           # 红
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light": RGBColor(0xD0, 0xE8, 0xE4),
        "gray": RGBColor(0x88, 0xAA, 0xA0),
        "table_h": RGBColor(0x00, 0x99, 0x80),
        "table_r1": RGBColor(0x14, 0x28, 0x24),
        "table_r2": RGBColor(0x1A, 0x32, 0x2E),
        "bg_overlay": "99",
        "bg_img_keywords": "medical health clean minimalist blue green",
    },
    "culture": {
        "name": "人文风",
        "scene": "文化/教育/历史",
        "bg_dark": RGBColor(0x1A, 0x14, 0x0E),
        "bg_dark2": RGBColor(0x25, 0x1E, 0x16),
        "card_bg": RGBColor(0x30, 0x28, 0x1E),
        "card_bg2": RGBColor(0x3A, 0x32, 0x26),
        "primary": RGBColor(0xC0, 0x60, 0x30),      # 赭石红
        "secondary": RGBColor(0xE8, 0xD8, 0xB0),     # 宣纸黄
        "accent": RGBColor(0xD0, 0xA0, 0x60),        # 古铜
        "highlight": RGBColor(0xA0, 0x50, 0x30),     # 深赭
        "purple": RGBColor(0x90, 0x60, 0x80),        # 紫
        "blue": RGBColor(0x60, 0x80, 0xA0),          # 灰蓝
        "red": RGBColor(0xC0, 0x40, 0x30),           # 红
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light": RGBColor(0xE0, 0xD8, 0xC8),
        "gray": RGBColor(0xA0, 0x90, 0x80),
        "table_h": RGBColor(0xA0, 0x60, 0x30),
        "table_r1": RGBColor(0x30, 0x28, 0x1E),
        "table_r2": RGBColor(0x3A, 0x32, 0x26),
        "bg_overlay": "99",
        "bg_img_keywords": "culture art vintage calligraphy warm paper",
    },
    "study": {
        "name": "学习风",
        "scene": "教育/培训/知识",
        "bg_dark": RGBColor(0x0E, 0x14, 0x1E),
        "bg_dark2": RGBColor(0x16, 0x1E, 0x2E),
        "card_bg": RGBColor(0x1E, 0x28, 0x3C),
        "card_bg2": RGBColor(0x26, 0x32, 0x4A),
        "primary": RGBColor(0x30, 0x80, 0xF0),      # 活力蓝
        "secondary": RGBColor(0xF0, 0xC0, 0x40),     # 明黄
        "accent": RGBColor(0x40, 0xE0, 0x80),        # 翠绿
        "highlight": RGBColor(0xF0, 0x60, 0x40),     # 橙红
        "purple": RGBColor(0x80, 0x60, 0xE0),        # 紫
        "blue": RGBColor(0x30, 0x80, 0xF0),          # 蓝
        "red": RGBColor(0xE0, 0x40, 0x50),           # 红
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light": RGBColor(0xD0, 0xE0, 0xF0),
        "gray": RGBColor(0x88, 0x99, 0xAA),
        "table_h": RGBColor(0x30, 0x70, 0xC0),
        "table_r1": RGBColor(0x1E, 0x28, 0x3C),
        "table_r2": RGBColor(0x26, 0x32, 0x4A),
        "bg_overlay": "99",
        "bg_img_keywords": "study education learning bright clean blue",
    },
}

# ============================================================
# 布局模板
# ============================================================
TEMPLATES = {
    "对比报告": {
        "description": "双栏对比 + 表格 + 卡片，适合产品/平台对比分析",
        "pages": ["封面", "背景", "章节过渡", "表格", "对比页", "对比页", "对比页", "章节过渡", "对比页", "三栏对比", "大表格", "优势卡片", "总结"],
    },
    "项目汇报": {
        "description": "项目进度/成果汇报",
        "pages": ["封面", "背景", "章节过渡", "卡片页", "表格", "卡片页", "总结"],
    },
    "产品介绍": {
        "description": "产品功能/卖点展示",
        "pages": ["封面", "背景", "章节过渡", "卡片页", "对比页", "卡片页", "总结"],
    },
}


def _c(style, key):
    """获取配色值"""
    return STYLES.get(style, STYLES["tech"])[key]


def _download_image(keywords, idx=0):
    """从网上搜索并下载配图（使用 Pollinations 或 placeholder）"""
    try:
        # 使用 pollinations.ai 生成配图（免费，无需 API Key）
        safe_kw = keywords.replace(" ", "+")
        url = f"https://image.pollinations.ai/prompt/{safe_kw}?width=1024&height=768&seed={random.randint(1, 99999)}"
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        fname = f"ppt_img_{ts}_{idx}_{random.randint(100,999)}.jpg"
        fpath = str(TEMP_DIR / fname)
        urllib.request.urlretrieve(url, fpath)
        if os.path.getsize(fpath) > 1000:
            return fpath
    except Exception:
        pass
    return None


def _add_bg_overlay(slide, w, h, alpha="99", color=None):
    """添加半透明遮罩"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color or RGBColor(0x08, 0x0C, 0x18)
    shape.line.fill.background()
    try:
        spPr = shape._element.find(qn('p:spPr'))
        if spPr is None:
            spPr = shape._element.find(qn('a:spPr'))
        if spPr is not None:
            solidFill = spPr.find(qn('a:solidFill'))
            if solidFill is not None:
                srgbClr = solidFill.find(qn('a:srgbClr'))
                if srgbClr is not None:
                    alpha_el = etree.SubElement(srgbClr, qn('a:alpha'))
                    alpha_el.set('val', alpha)
    except Exception:
        pass
    return shape


def _add_bg_img(slide, w, h, img_path=None, style="tech"):
    """添加背景图+遮罩"""
    path = img_path
    if path and os.path.exists(path):
        try:
            slide.shapes.add_picture(path, Inches(0), Inches(0), w, h)
        except Exception:
            pass
    _add_bg_overlay(slide, w, h, _c(style, "bg_overlay"), _c(style, "bg_dark"))


def _add_dark_bg(slide, w, h, style="tech"):
    """纯色深色背景"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _c(style, "bg_dark")
    shape.line.fill.background()
    # 顶部装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), w, Pt(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _c(style, "primary")
    bar.line.fill.background()


def _add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def _add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=18, color=None, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    if color is None:
        color = RGBColor(0xFF, 0xFF, 0xFF)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def _add_multiline_text(slide, left, top, width, height, lines, font_size=14, color=None, line_spacing=1.5, font_name='Microsoft YaHei'):
    if color is None:
        color = RGBColor(0xD8, 0xE0, 0xEE)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.3)
    return txBox


def _add_icon_circle(slide, left, top, size, text, fill_color, font_size=20):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    shape.text_frame.margin_left = Pt(0)
    shape.text_frame.margin_right = Pt(0)
    shape.text_frame.margin_top = Pt(0)
    shape.text_frame.margin_bottom = Pt(0)
    return shape


def _add_accent_line(slide, left, top, width, height_pt=4, color=None):
    if color is None:
        color = RGBColor(0x00, 0xD4, 0xE8)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(height_pt))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_table(slide, left, top, width, height, data, col_widths=None, font_size=12, style="tech"):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r_idx, row_data in enumerate(data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = 'Microsoft YaHei'
                paragraph.alignment = PP_ALIGN.CENTER
                if r_idx == 0:
                    paragraph.font.color.rgb = _c(style, "white")
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = _c(style, "light")
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _c(style, "table_h")
            elif r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _c(style, "table_r1")
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _c(style, "table_r2")
    return table_shape


def _add_image(slide, path, left, top, width, height=None):
    if path and os.path.exists(path):
        try:
            if height:
                slide.shapes.add_picture(path, left, top, width, height)
            else:
                slide.shapes.add_picture(path, left, top, width)
            return True
        except Exception:
            pass
    return False


def _add_comparison_card(slide, x, y, w, h, icon, icon_color, title, items, border_color, font_size=22, style="tech"):
    """对比卡片"""
    _add_rounded_rect(slide, x, y, w, h, _c(style, "card_bg"), border_color)
    _add_icon_circle(slide, x + Inches(0.25), y + Inches(0.2), Inches(0.6), icon, icon_color, 20)
    _add_text_box(slide, x + Inches(1.0), y + Inches(0.25), w - Inches(1.3), Inches(0.45),
                  title, 24, icon_color, True)
    _add_multiline_text(slide, x + Inches(0.3), y + Inches(0.9), w - Inches(0.6), h - Inches(1.1),
                        items, font_size, _c(style, "light"))


def _add_img_card(slide, x, y, w, h, img_path, label, border_color, style="tech"):
    """带图片的卡片"""
    _add_rounded_rect(slide, x, y, w, h, _c(style, "card_bg2"), border_color)
    if _add_image(slide, img_path, x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), h - Inches(0.9)):
        _add_text_box(slide, x + Inches(0.15), y + h - Inches(0.65), w - Inches(0.3), Inches(0.5),
                      label, 11, _c(style, "gray"), align=PP_ALIGN.CENTER)
        return True
    return False


# ============================================================
# 页面生成函数
# ============================================================

def _make_cover(prs, w, h, title, subtitle, style, bg_img):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_img(slide, w, h, bg_img, style)

    _add_rect(slide, Inches(0), Inches(0), w, Inches(0.06), _c(style, "primary"))
    _add_rect(slide, Inches(0.6), Inches(1.5), Inches(0.06), Inches(4.5), _c(style, "secondary"))
    _add_text_box(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(1.2),
                  title, 44, _c(style, "white"), True)
    _add_text_box(slide, Inches(1.2), Inches(3.2), Inches(10), Inches(0.8),
                  subtitle or "", 30, _c(style, "primary"))
    _add_accent_line(slide, Inches(1.2), Inches(4.2), Inches(3.5), 4, _c(style, "accent"))
    _add_text_box(slide, Inches(9.5), Inches(6.8), Inches(3.5), Inches(0.5),
                  f"{datetime.now().strftime('%Y年%m月')}  |  坐山客研究团队", 14, _c(style, "gray"), align=PP_ALIGN.RIGHT)


def _make_section_transition(prs, w, h, number, title, style, bg_img):
    """章节过渡页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_img(slide, w, h, bg_img, style)

    _add_rect(slide, Inches(0), Inches(3.2), Inches(0.12), Inches(1.2), _c(style, "primary"))
    _add_text_box(slide, Inches(0.6), Inches(2.8), Inches(4), Inches(0.8),
                  number, 60, _c(style, "primary"), True)
    _add_text_box(slide, Inches(0.6), Inches(3.6), Inches(8), Inches(0.8),
                  title, 38, _c(style, "white"), True)
    _add_accent_line(slide, Inches(0.6), Inches(4.5), Inches(3.5), 3, _c(style, "accent"))


def _make_background_page(prs, w, h, title, cards_data, style, bg_img=None):
    """背景介绍页（4卡片）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_dark_bg(slide, w, h, style)

    _add_text_box(slide, Inches(0.8), Inches(0.2), Inches(6), Inches(0.5),
                  title, 34, _c(style, "white"), True)
    _add_accent_line(slide, Inches(0.8), Inches(0.7), Inches(2.5), 3, _c(style, "primary"))

    # 左侧配图
    if bg_img:
        _add_img_card(slide, Inches(0.3), Inches(1.1), Inches(5.0), Inches(5.8),
                      bg_img, "▲ 配图", _c(style, "primary"), style)

    card_x = Inches(5.8) if bg_img else Inches(0.3)
    card_w = Inches(7.0) if bg_img else Inches(12.5)

    for i, card in enumerate(cards_data):
        if isinstance(card, dict):
            icon = card.get("icon", "📌")
            color_key = card.get("color", "primary")
            ctitle = card.get("title", "")
            items = card.get("items", [])
        else:
            icon, color_key, ctitle, items = card
        x = card_x
        y = Inches(1.1 + i * 1.5)
        hh = Inches(1.3)
        ww = card_w
        color = _c(style, color_key)
        _add_rounded_rect(slide, x, y, ww, hh, _c(style, "card_bg2"), color)
        _add_icon_circle(slide, x + Inches(0.15), y + Inches(0.25), Inches(0.7), icon, color, 22)
        _add_text_box(slide, x + Inches(1.0), y + Inches(0.15), ww - Inches(1.3), Inches(0.4),
                      ctitle, 20, color, True)
        _add_multiline_text(slide, x + Inches(1.0), y + Inches(0.55), ww - Inches(1.3), hh - Inches(0.6),
                            items, 14, _c(style, "light"))


def _make_comparison_page(prs, w, h, title, left_data, right_data, style, img_path=None):
    """双栏对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_dark_bg(slide, w, h, style)

    _add_text_box(slide, Inches(0.8), Inches(0.12), Inches(10), Inches(0.45),
                  title, 30, _c(style, "white"), True)
    _add_accent_line(slide, Inches(0.8), Inches(0.55), Inches(2.5), 3, _c(style, "primary"))

    # 顶部配图区
    if img_path:
        _add_img_card(slide, Inches(0.3), Inches(0.8), Inches(5.5), Inches(2.8),
                      img_path, "▲ 截图", _c(style, "blue"), style)

        _add_icon_circle(slide, Inches(5.9), Inches(1.6), Inches(1.4), "VS", _c(style, "accent"), 22)

        _add_rounded_rect(slide, Inches(7.5), Inches(0.8), Inches(5.5), Inches(2.8),
                          _c(style, "card_bg2"), _c(style, "secondary"))
        _add_text_box(slide, Inches(7.8), Inches(1.0), Inches(5.0), Inches(0.5),
                      right_data.get("title", ""), 22, _c(style, "secondary"), True)
        _add_multiline_text(slide, Inches(7.8), Inches(1.6), Inches(5.0), Inches(1.8),
                            right_data.get("desc", []), 16, _c(style, "light"))

        top_offset = Inches(3.9)
    else:
        top_offset = Inches(0.8)

    card_h = H - top_offset - Inches(0.3)

    _add_comparison_card(slide, Inches(0.3), top_offset, Inches(5.8), card_h,
                         left_data.get("icon", "L"), _c(style, left_data.get("color", "blue")),
                         left_data.get("title", ""), left_data.get("items", []),
                         _c(style, left_data.get("color", "blue")), font_size=22, style=style)

    _add_comparison_card(slide, Inches(7.0), top_offset, Inches(5.8), card_h,
                         right_data.get("icon", "R"), _c(style, right_data.get("color", "secondary")),
                         right_data.get("title", ""), right_data.get("items", []),
                         _c(style, right_data.get("color", "secondary")), font_size=22, style=style)


def _make_three_column_page(prs, w, h, title, columns, style):
    """三栏对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_dark_bg(slide, w, h, style)

    _add_text_box(slide, Inches(0.8), Inches(0.12), Inches(10), Inches(0.45),
                  title, 28, _c(style, "white"), True)
    _add_accent_line(slide, Inches(0.8), Inches(0.55), Inches(2.5), 3, _c(style, "primary"))

    card_w3 = Inches(4.0)
    card_h3 = Inches(5.8)

    for i, col in enumerate(columns):
        x = Inches(0.3 + i * 4.3)
        _add_comparison_card(slide, x, Inches(0.9), card_w3, card_h3,
                             col.get("icon", "?"), _c(style, col.get("color", "blue")),
                             col.get("title", ""), col.get("items", []),
                             _c(style, col.get("color", "blue")), font_size=20, style=style)


def _make_table_page(prs, w, h, title, table_data, col_widths, style):
    """表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_dark_bg(slide, w, h, style)

    _add_text_box(slide, Inches(0.8), Inches(0.12), Inches(8), Inches(0.45),
                  title, 30, _c(style, "white"), True)
    _add_accent_line(slide, Inches(0.8), Inches(0.55), Inches(2.5), 3, _c(style, "primary"))

    _add_table(slide, Inches(0.8), Inches(0.8), Inches(10.0), Inches(6.3),
               table_data, col_widths, font_size=14, style=style)


def _make_card_grid_page(prs, w, h, title, cards, cols=4, style="tech"):
    """卡片网格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_dark_bg(slide, w, h, style)

    _add_text_box(slide, Inches(0.8), Inches(0.15), Inches(8), Inches(0.45),
                  title, 30, _c(style, "white"), True)
    _add_accent_line(slide, Inches(0.8), Inches(0.55), Inches(2.5), 3, _c(style, "primary"))

    for i, card in enumerate(cards):
        row = i // cols
        col = i % cols
        x = Inches(0.5 + col * 3.1)
        y = Inches(0.9 + row * 3.2)
        cw = Inches(2.8)
        ch = Inches(2.9)
        color = _c(style, card.get("color", "primary"))
        _add_rounded_rect(slide, x, y, cw, ch, _c(style, "card_bg2"), color)
        _add_icon_circle(slide, x + Inches(0.9), y + Inches(0.15), Inches(0.55),
                         card.get("icon", "?"), color, 16)
        _add_text_box(slide, x + Inches(0.1), y + Inches(0.75), cw - Inches(0.2), Inches(0.35),
                      card.get("title", ""), 16, color, True, PP_ALIGN.CENTER)
        _add_accent_line(slide, x + Inches(0.4), y + Inches(1.1), Inches(2.0), 2, color)
        _add_multiline_text(slide, x + Inches(0.1), y + Inches(1.3), cw - Inches(0.2), ch - Inches(1.5),
                            card.get("items", []), 13, _c(style, "light"), 1.3)


def _make_summary_page(prs, w, h, title, points, style):
    """总结页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_dark_bg(slide, w, h, style)

    _add_rect(slide, Inches(0), Inches(0), w, Inches(0.06), _c(style, "primary"))
    _add_rect(slide, Inches(0.6), Inches(1.5), Inches(0.06), Inches(4.5), _c(style, "accent"))
    _add_text_box(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(0.8),
                  title, 38, _c(style, "white"), True)
    _add_accent_line(slide, Inches(1.2), Inches(2.3), Inches(3), 3, _c(style, "primary"))

    _add_multiline_text(slide, Inches(1.2), Inches(2.8), Inches(10.5), Inches(3.5),
                        points, 20, _c(style, "light"), 1.8)

    _add_accent_line(slide, Inches(1.2), Inches(6.3), Inches(5), 2, _c(style, "accent"))
    if points:
        _add_text_box(slide, Inches(1.2), Inches(6.5), Inches(10), Inches(0.6),
                      points[-1] if isinstance(points[-1], str) else "",
                      18, _c(style, "accent"))


# ============================================================
# 主入口
# ============================================================

def gen_ppt_pro(
    title: str = "演示文稿",
    subtitle: str = "",
    slides: list = None,
    style: str = "tech",
    template: str = "对比报告",
    author: str = "坐山客",
    output_name: str = "",
    bg_img: str = None,
) -> str:
    """
    生成专业版 PowerPoint (.pptx) 演示文稿

    Args:
        title: 演示文稿标题
        subtitle: 副标题
        slides: 幻灯片数据列表。每项格式:
            {
                "type": "cover|transition|background|comparison|three_column|table|card_grid|summary",
                "title": "页面标题",
                ... 其他字段取决于 type
            }
        style: 风格配色 (tech/fashion/food/medical/culture/study)
        template: 模板名称（"对比报告"/"项目汇报"/"产品介绍"）
        author: 作者
        output_name: 输出文件名（不含 .pptx）
        bg_img: 背景图路径（可选）

    Returns:
        JSON: {success, file_path, slide_count, title, style, error?}
    """
    if Presentation is None:
        return json.dumps({"success": False, "error": "python-pptx 未安装，请运行: pip install python-pptx"})

    try:
        if slides is None:
            slides = []

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        W = prs.slide_width
        H = prs.slide_height

        # 下载背景图（如果没有提供）
        if not bg_img:
            bg_img = _download_image(_c(style, "bg_img_keywords"), 0)

        for sec in slides:
            stype = sec.get("type", "cover")

            if stype == "cover":
                _make_cover(prs, W, H, sec.get("title", title),
                           sec.get("subtitle", subtitle), style, bg_img)

            elif stype == "transition":
                _make_section_transition(prs, W, H, sec.get("number", "01"),
                                        sec.get("title", ""), style, bg_img)

            elif stype == "background":
                _make_background_page(prs, W, H, sec.get("title", "报告背景"),
                                     sec.get("cards", []), style, bg_img)

            elif stype == "comparison":
                _make_comparison_page(prs, W, H, sec.get("title", ""),
                                     sec.get("left", {}), sec.get("right", {}),
                                     style, sec.get("img_path"))

            elif stype == "three_column":
                _make_three_column_page(prs, W, H, sec.get("title", ""),
                                       sec.get("columns", []), style)

            elif stype == "table":
                _make_table_page(prs, W, H, sec.get("title", ""),
                                sec.get("data", []), sec.get("col_widths"), style)

            elif stype == "card_grid":
                _make_card_grid_page(prs, W, H, sec.get("title", ""),
                                    sec.get("cards", []), sec.get("cols", 4), style)

            elif stype == "summary":
                _make_summary_page(prs, W, H, sec.get("title", "最终结论"),
                                  sec.get("points", []), style)

        # 保存
        if not output_name:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_title = re.sub(r"[^\\w\\u4e00-\\u9fff_-]", "_", title)[:20]
            output_name = f"{safe_title}_{style}_{ts}"
        output_path = str(OUTPUT_DIR / f"{output_name}.pptx")
        prs.save(output_path)

        return json.dumps({
            "success": True,
            "file_path": output_path,
            "slide_count": len(slides),
            "title": title,
            "style": style,
            "style_name": STYLES.get(style, {}).get("name", "未知"),
        }, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"success": False, "error": str(e)}, ensure_ascii=False)
